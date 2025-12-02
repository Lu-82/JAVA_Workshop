from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide content as (title, content)
slides_content = [
    ("Chain Surveying", "Basics, Instruments, Methods & Applications\nPresented by: [Your Name]\nDate: [Presentation Date]"),
    ("Introduction", "• Chain surveying is the simplest method of surveying.\n• It is based on linear measurements only.\n• Used when the area is small, flat, and free of obstacles."),
    ("Objectives", "• Understand the principles of chain surveying.\n• Learn about the instruments used.\n• Explore the procedure and types of lines.\n• Know the advantages and limitations."),
    ("Instruments Used", "1. Chain (Metric chain – 20m/30m)\n2. Tape (for short measurements)\n3. Arrows (marking end of chain length)\n4. Ranging Rods (for alignment)\n5. Peg (marking stations)\n6. Plumb Bob (used in sloped ground)\n7. Cross Staff (for perpendicular lines)"),
    ("Principle of Chain Surveying", "• Based on the concept of triangulation.\n• Main lines and check lines form triangles.\n• Accurate plotting using measured distances only."),
    ("Types of Lines in Chain Surveying", "1. Base Line – main and longest line\n2. Chain Lines – straight lines between stations\n3. Tie Lines – connect interior details to chain lines\n4. Check Lines – used to check accuracy"),
    ("Procedure", "1. Reconnaissance\n2. Station Marking\n3. Ranging\n4. Chaining (measurement)\n5. Recording in field book\n6. Plotting"),
    ("Advantages", "• Simple and inexpensive\n• Requires few instruments\n• Good for small, flat areas"),
    ("Limitations", "• Not suitable for hilly areas\n• Less accurate over long distances\n• Cannot record elevation differences"),
    ("Applications", "• Plotting small plots or fields\n• Road and canal surveys\n• Boundary surveys"),
    ("Conclusion", "• Chain surveying is a fundamental method in civil engineering.\n• Understanding its principles is crucial for land measurement and mapping."),
    ("Thank You", "Any Questions?")
]

# Create slides
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Save the presentation
prs.save("Chain_Surveying_Presentation.pptx")